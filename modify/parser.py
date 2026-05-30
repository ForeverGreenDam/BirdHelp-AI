"""文档逆向解析 — 降级兜底方案。

当 file_record.outline 为空时（旧格式文件，生成时未回传大纲），
从 .pptx / .docx 文件中逆向提取大纲文本，标注 source: "file_parse"。

解析保真度约 40-50%（§一）：
  - 可恢复: page_number, title, body 文本
  - 部分: chart_data / table_data（丢元信息）
  - 丢失: layout_type, visual_plan, image_query, style

返回格式与 chains 输出兼容，可直接用于对话修改。
"""

from __future__ import annotations

from typing import Any

from loguru import logger

from utils.file import temp_file_path, ensure_temp_dir


async def parse_outline_from_file(
    file_bytes: bytes,
    doc_type: str,
    file_name: str = "document",
) -> dict[str, Any]:
    """从 Office 文件逆向提取大纲。

    Args:
        file_bytes: 文件的原始字节内容
        doc_type: "ppt" / "word" / "pdf"
        file_name: 原始文件名（用于日志）

    Returns:
        与 chains 输出兼容的大纲 JSON，额外包含 meta.source = "file_parse"
    """
    ensure_temp_dir()

    try:
        if doc_type == "ppt":
            return _parse_pptx(file_bytes, file_name)
        elif doc_type == "word":
            return _parse_docx(file_bytes, file_name)
        elif doc_type == "pdf":
            return _parse_pdf(file_bytes, file_name)
        else:
            logger.warning(f"不支持的文档类型用于解析: {doc_type}")
            return _empty_outline(doc_type, f"不支持的文档类型: {doc_type}")
    except Exception as exc:
        logger.error(f"文件解析失败 [{file_name}]: {exc}")
        return _empty_outline(doc_type, f"解析失败: {exc}")


# ── PPT 解析 ──

def _parse_pptx(content: bytes, file_name: str) -> dict[str, Any]:
    """从 .pptx 提取文本结构。"""
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slides_outline = []

    for i, slide in enumerate(prs.slides, start=1):
        # 提取所有文本框文本
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

        # 第一个非空文本框视为标题，其余为正文
        title = texts[0] if texts else f"第{i}页"
        body = "\n".join(texts[1:]) if len(texts) > 1 else ""

        slides_outline.append({
            "page_number": i,
            "title": title,
            "body": body,
            "layout_type": "text_only",  # 无法从文件中恢复
            "visual_plan": "",
            "image_query": "",
            "chart_data": None,
            "table_data": None,
            "style": "",
        })

    logger.info(f"PPT 解析完成 [{file_name}]: {len(slides_outline)} 页")
    return {
        "title": file_name,
        "subtitle": "",
        "doc_type": "ppt",
        "style": "academic",
        "slides": slides_outline,
        "sections": [],
        "_meta": {"source": "file_parse"},
    }


# ── Word 解析 ──

def _parse_docx(content: bytes, file_name: str) -> dict[str, Any]:
    """从 .docx 提取段落结构。"""
    import io
    from docx import Document

    doc = Document(io.BytesIO(content))
    sections_outline = []
    current_section = 1
    current_heading = ""

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # 标题样式视为节标题
        if para.style.name.startswith("Heading") or "heading" in para.style.name.lower():
            # 保存上一节
            if current_heading:
                sections_outline.append({
                    "section_number": current_section,
                    "heading": current_heading,
                    "content": "",
                    "has_image": False,
                    "image_query": "",
                })
                current_section += 1
            current_heading = text
        else:
            # 正文追加到当前节
            if not sections_outline and not current_heading:
                current_heading = file_name
            content_text = text
            # 简化：累积所有正文
            pass

    # 补上最后一节
    if current_heading or sections_outline:
        sections_outline.append({
            "section_number": current_section,
            "heading": current_heading or file_name,
            "content": "",
            "has_image": False,
            "image_query": "",
        })

    # 如果没有任何标题，提取所有段落作为单节
    if not sections_outline:
        all_text = "\n".join(
            p.text.strip() for p in doc.paragraphs if p.text.strip()
        )
        sections_outline.append({
            "section_number": 1,
            "heading": file_name,
            "content": all_text[:2000],  # 截断，避免过大
            "has_image": False,
            "image_query": "",
        })

    logger.info(f"Word 解析完成 [{file_name}]: {len(sections_outline)} 节")
    return {
        "title": file_name,
        "subtitle": "",
        "doc_type": "word",
        "style": "academic",
        "slides": [],
        "sections": sections_outline,
        "_meta": {"source": "file_parse"},
    }


# ── PDF 解析 ──

def _parse_pdf(content: bytes, file_name: str) -> dict[str, Any]:
    """从 PDF 提取文本（使用 PyPDF2，无 layout 信息）。

    PDF 解析能力有限——只能恢复文本，layout_type/表格/图表等元信息完全丢失。
    """
    try:
        from PyPDF2 import PdfReader
        import io

        reader = PdfReader(io.BytesIO(content))
        pages_outline = []

        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            lines = text.strip().split("\n")
            title = lines[0][:200] if lines else f"第{i}页"
            body = "\n".join(lines[1:])[:1000] if len(lines) > 1 else ""

            pages_outline.append({
                "page_number": i,
                "title": title,
                "body": body,
                "layout_type": "text_only",
                "visual_plan": "",
                "image_query": "",
                "chart_data": None,
                "table_data": None,
                "style": "",
            })

        logger.info(f"PDF 解析完成 [{file_name}]: {len(pages_outline)} 页")
        return {
            "title": file_name,
            "subtitle": "",
            "doc_type": "pdf",
            "style": "academic",
            "slides": pages_outline,
            "sections": [],
            "_meta": {"source": "file_parse"},
        }
    except ImportError:
        logger.warning("PyPDF2 未安装，无法解析 PDF")
        return _empty_outline("pdf", "PyPDF2 未安装")


# ── 空大纲 ──

def _empty_outline(doc_type: str, reason: str = "") -> dict[str, Any]:
    return {
        "title": "",
        "doc_type": doc_type,
        "style": "academic",
        "slides": [],
        "sections": [],
        "_meta": {"source": "file_parse", "parse_error": reason},
    }

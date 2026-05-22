"""PPT 文件生成器 — 将 LLM 输出的结构化视觉描述构建为 .pptx 文件。

Generator 作为编排入口，遍历 slides 并根据每页的 layout_type
分发到对应的布局渲染器（cover / section / text_only / text_image 等）。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from loguru import logger

from generator.base import BaseGenerator
from generator.ppt.layout import create_dna
from generator.ppt.layouts import dispatch_renderer
from generator.ppt.shapes import SLIDE_W, SLIDE_H

# 统一使用 blank 布局（索引 6），所有渲染器自行绘制全部元素
_BLANK_LAYOUT_IDX = 6


class PptGenerator(BaseGenerator):
    """PPT 生成器 — 使用设计系统 + 布局渲染器生成视觉丰富的幻灯片。

    每页根据 layout_type 字段分发给对应的渲染器，
    渲染器接收 theme/dna/images 参数完成绘制。
    """

    output_extension = ".pptx"

    def generate(
        self,
        content: dict[str, Any],
        output_path: Path,
        images_map: dict[str, list[str]] | None = None,
    ) -> Path:
        """根据结构化视觉描述生成 PPT 文件。

        Args:
            content: LLM 输出的结构化 JSON（含 layout_type/visual_plan/image_query）
            output_path: 输出 .pptx 文件路径
            images_map: slide_key → 本地图片路径列表

        Returns:
            生成的文件路径
        """
        parsed = self._parse_content(content)
        slides = parsed.get("slides", [])
        style_name = parsed.get("style", "academic")
        images_map = images_map or {}

        dna = create_dna(style_name, parsed.get("title", ""))
        theme = dna.theme

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        total = len(slides)
        for i, slide_data in enumerate(slides):
            slide_data["_total_pages"] = total
            slide_data.setdefault("page_number", i + 1)

            page_num = slide_data.get("page_number", i + 1)
            img_key = f"slide_{page_num:02d}"
            imgs = images_map.get(img_key, [])

            # 兜底：图片映射中没有对应图片时，为该页当场生成占位图
            if not imgs and slide_data.get("image_query", "").strip():
                from generator.ppt.image_provider import _generate_placeholder, _images_dir, _query_hash
                query = slide_data["image_query"].strip()
                placeholder_path = _images_dir() / f"{img_key}-fallback.png"
                if _generate_placeholder(query, placeholder_path):
                    imgs = [str(placeholder_path)]

            layout_type = slide_data.get("layout_type", "text_only")

            slide_layout = prs.slide_layouts[_BLANK_LAYOUT_IDX]
            slide = prs.slides.add_slide(slide_layout)

            renderer = dispatch_renderer(layout_type)
            renderer(slide, slide_data, theme, dna, imgs)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info(
            f"PPT generated: {output_path}, {total} slides, "
            f"style={style_name}, dna={dna.shape_style}/{dna.density}/{dna.decoration_level}"
        )
        return output_path

"""形状工具包 — 封装 python-pptx 常用形状操作，提供声明式绘图接口。

上层布局渲染器通过本模块的函数描述"画什么"（颜色、位置、尺寸），
无需直接操作 python-pptx 的 XML 级别 API。
所有尺寸参数使用 Inches() 等 EMU 值。

Phase 2 新增:
- 富文本支持（粗体/斜体/超链接标记解析）
- 菱形/箭头/星形等增强形状
"""

from __future__ import annotations

import re
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from generator.ppt.theme import ColorTheme

# ── 富文本标记正则 ──
# **粗体**  *斜体*  [文本](url)
_BOLD_RE = re.compile(r'\*\*(.+?)\*\*')
_ITALIC_RE = re.compile(r'\*(.+?)\*')
_LINK_RE = re.compile(r'\[(.+?)\]\((.+?)\)')

# ── 幻灯片尺寸（16:9） ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SAFE_L = Inches(0.8)   # 左边距
SAFE_R = Inches(12.533)  # 右边距 = slide_w - 0.8
SAFE_T = Inches(0.4)   # 上边距
SAFE_B = Inches(7.1)   # 下边距 = slide_h - 0.4


def add_rect(
    slide,
    left, top, width, height,
    fill_color: RGBColor | str,
    border_color: RGBColor | str | None = None,
    border_width: float = 0,
    corner_radius: int = 0,
) -> object:
    """添加矩形色块。corner_radius > 0 时使用圆角矩形。"""
    if corner_radius > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        # 调整圆角半径（通过 XML 属性）
        try:
            prst_geom = shape._element.find(qn('a:prstGeom'))
            if prst_geom is not None:
                avlst = prst_geom.find(qn('a:avLst'))
                if avlst is not None:
                    gd = avlst.find(qn('a:gd'))
                    if gd is not None:
                        gd.set('fmla', f'val {corner_radius * 12700}')  # EMU: pt * 12700
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
    # 填充
    shape.fill.solid()
    if isinstance(fill_color, str):
        fill_color = _hex_to_rgb(fill_color)
    shape.fill.fore_color.rgb = fill_color
    # 边框
    if border_width > 0 and border_color:
        shape.line.color.rgb = border_color if isinstance(border_color, RGBColor) else _hex_to_rgb(border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_circle(
    slide,
    cx: Inches, cy: Inches, radius: Inches,
    fill_color: RGBColor | str,
) -> object:
    """添加圆形装饰元素。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - radius, cy - radius,
        radius * 2, radius * 2,
    )
    shape.fill.solid()
    if isinstance(fill_color, str):
        fill_color = _hex_to_rgb(fill_color)
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_accent_bar(
    slide,
    left, top, width, height,
    color: RGBColor | str,
) -> object:
    """添加侧边强调条（细长矩形）。"""
    return add_rect(slide, left, top, width, height, color, border_width=0)


def add_line(
    slide,
    x1, y1, x2, y2,
    color: RGBColor | str,
    width: float = 1.5,
) -> object:
    """添加直线。"""
    if isinstance(color, str):
        color = _hex_to_rgb(color)
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT = 1
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_text_box(
    slide,
    left, top, width, height,
    text: str,
    font_name: str,
    font_size: int,
    font_color: RGBColor | str,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,  # noqa: ANN001 (PP_ALIGN is IntEnum)
    word_wrap: bool = True,
    vertical_anchor=MSO_ANCHOR.TOP,  # noqa: ANN001
) -> object:
    """添加统一样式的文本框，返回 shape 对象以便调用方进一步操作。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    # 设置垂直锚定
    try:
        txBox.text_frame._txBody.bodyPr.set('anchor', {
            MSO_ANCHOR.TOP: 't',
            MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b',
        }.get(vertical_anchor, 't'))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if isinstance(font_color, str):
        font_color = _hex_to_rgb(font_color)
    p.font.color.rgb = font_color
    return txBox


def add_multiline_text_box(
    slide,
    left, top, width, height,
    lines: list[dict],
    theme: ColorTheme,
) -> object:
    """添加多段落文本框，每行可独立设置样式。

    lines 格式: [{"text": "...", "size": 18, "bold": False, "color": theme.text_body, "bullet": True}, ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "• " if line.get("bullet", False) else ""
        p.text = f"{prefix}{line['text']}"
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.font.name = theme.body_font
        p.font.size = Pt(line.get("size", 18))
        p.font.bold = line.get("bold", False)
        color = line.get("color", theme.text_body)
        if isinstance(color, str):
            color = _hex_to_rgb(color)
        p.font.color.rgb = color
        if line.get("space_after"):
            p.space_after = Pt(line["space_after"])
        if line.get("space_before"):
            p.space_before = Pt(line["space_before"])
    return txBox


def add_image(
    slide,
    left, top, width, height,
    image_path: str,
) -> object:
    """插入图片，自动等比缩放填充指定区域（居中裁剪）。"""
    from io import BytesIO
    from PIL import Image
    from loguru import logger

    # 先尝试直接用 python-pptx 原生插入（兼容性最好）
    try:
        return slide.shapes.add_picture(image_path, left, top, width, height)
    except Exception as native_err:
        logger.debug(f"add_picture native failed for {image_path}: {native_err}, trying PIL crop")

    try:
        with Image.open(image_path) as img:
            img_w, img_h = img.size
    except Exception as pil_err:
        logger.warning(f"add_image: both native and PIL failed for {image_path}: {pil_err}")
        raise

    slot_ratio = width / height
    img_ratio = img_w / img_h

    if abs(img_ratio - slot_ratio) < 0.01:
        buf = BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)
        return slide.shapes.add_picture(buf, left, top, width, height)

    # 等比缩放后居中裁剪
    if img_ratio > slot_ratio:
        new_w = int(img_h * slot_ratio)
        offset_x = (img_w - new_w) // 2
        cropped = img.crop((offset_x, 0, offset_x + new_w, img_h))
    else:
        new_h = int(img_w / slot_ratio)
        offset_y = (img_h - new_h) // 2
        cropped = img.crop((0, offset_y, img_w, offset_y + new_h))

    buf = BytesIO()
    cropped.save(buf, format='PNG')
    buf.seek(0)
    return slide.shapes.add_picture(buf, left, top, width, height)


def set_slide_bg(slide, color: RGBColor | str) -> None:
    """设置幻灯片纯色背景。"""
    if isinstance(color, str):
        color = _hex_to_rgb(color)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_page_number(
    slide,
    index: int,
    total: int,
    theme: ColorTheme,
) -> None:
    """在右下角添加页码 'index / total'。"""
    from pptx.util import Inches
    page_text = f"{index} / {total}"
    add_text_box(
        slide,
        Inches(11.5), Inches(7.0),
        Inches(1.5), Inches(0.4),
        page_text,
        font_name=theme.body_font,
        font_size=10,
        font_color=theme.text_secondary,
        alignment=PP_ALIGN.RIGHT,
    )


def clear_placeholders(slide) -> None:
    """移除幻灯片的所有占位符，避免空占位符导致的兼容性警告。"""
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


# ── 富文本文本框 ──

def add_rich_text_box(
    slide,
    left, top, width, height,
    text: str,
    font_name: str,
    font_size: int,
    default_color: RGBColor | str,
    alignment=PP_ALIGN.LEFT,  # noqa: ANN001
    line_spacing: float = 1.3,
) -> object:
    """支持轻量标记的富文本文本框。

    支持的标记语法:
    - **粗体文字**
    - *斜体文字*
    - [链接文字](https://url)
    标记和链接均可嵌套于普通文字中。

    Args:
        slide: pptx slide 对象
        left, top, width, height: 文本框位置尺寸 (EMU)
        text: 含标记的文本
        font_name: 字体名
        font_size: 基准字号 (pt)
        default_color: 默认文字颜色
        alignment: 对齐方式
        line_spacing: 行间距倍数
    """
    if isinstance(default_color, str):
        default_color = _hex_to_rgb(default_color)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.alignment = alignment
    # 行间距通过 XML 属性设置
    try:
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    except Exception:
        pass

    _apply_rich_text_to_paragraph(p, text, font_name, font_size, default_color)

    return txBox


def _apply_rich_text_to_paragraph(paragraph, text: str, font_name: str, font_size: int, default_color) -> None:
    """将带标记的文本解析为多个 run，应用对应样式。"""
    # 使用正则将文本拆分为 (type, content) 序列
    # type: "text", "bold", "italic", "link"
    token_pattern = re.compile(
        r'(\*\*(?P<bold>.+?)\*\*)'
        r'|(\*(?P<italic>.+?)\*)'
        r'|(\[(?P<link_text>.+?)\]\((?P<link_url>.+?)\))'
    )
    segments = []  # list of (start, end, type, extra)
    for m in token_pattern.finditer(text):
        if m.group('bold'):
            segments.append((m.start(), m.end(), 'bold', m.group('bold')))
        elif m.group('italic'):
            segments.append((m.start(), m.end(), 'italic', m.group('italic')))
        elif m.group('link_text'):
            segments.append((m.start(), m.end(), 'link', (m.group('link_text'), m.group('link_url'))))

    if not segments:
        # 纯文本
        run = paragraph.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = default_color
        return

    # 在标记之间插入普通文本
    pos = 0
    for start, end, seg_type, content in segments:
        if pos < start:
            _add_run(paragraph, text[pos:start], font_name, font_size, default_color, False, False)
        if seg_type == 'bold':
            _add_run(paragraph, content, font_name, font_size, default_color, True, False)
        elif seg_type == 'italic':
            _add_run(paragraph, content, font_name, font_size, default_color, False, True)
        elif seg_type == 'link':
            link_text, link_url = content
            run = paragraph.add_run()
            run.text = link_text
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = default_color
            # 超链接通过 XML 设置
            try:
                rPr = run._r.get_or_add_rPr()
                hlinkClick = rPr.makeelement(qn('a:hlinkClick'), {
                    qn('r:id'): _add_hyperlink_relationship(paragraph, link_url),
                })
                rPr.insert(0, hlinkClick)
            except Exception:
                pass
        pos = end
    if pos < len(text):
        _add_run(paragraph, text[pos:], font_name, font_size, default_color, False, False)


def _add_run(paragraph, text: str, font_name: str, font_size: int, color, bold: bool, italic: bool):
    """向段落添加一个文字 run。"""
    if not text:
        return
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _add_hyperlink_relationship(paragraph, url: str) -> str:
    """为超链接添加关系，返回 rId。"""
    # 利用 paragraph 所属的 slide 来添加关系
    # 从 paragraph._p 向上查找 slide 元素
    el = paragraph._p
    while el is not None:
        parent = el.getparent()
        if parent is None:
            break
        # 查找包含关系的元素（通常是 slide.xml 的根元素下）
        el = parent
        # 尝试找到 rId 的起始点
    # 简化处理：通过 slide 的 part 来添加关系
    try:
        # 获取 slide part
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        slide_part = paragraph.part
        rId = slide_part.relate_to(url, RT.HYPERLINK, is_external=True)
        return rId
    except Exception:
        return ""


# ── 增强形状库 ──

def add_diamond(
    slide,
    left, top, width, height,
    fill_color: RGBColor | str,
    border_color: RGBColor | str | None = None,
    border_width: float = 0,
) -> object:
    """添加菱形装饰。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, width, height)
    _apply_fill_border(shape, fill_color, border_color, border_width)
    return shape


def add_arrow_shape(
    slide,
    left, top, width, height,
    fill_color: RGBColor | str,
    direction: str = "right",  # "right" | "left" | "up" | "down"
    border_color: RGBColor | str | None = None,
    border_width: float = 0,
) -> object:
    """添加箭头形状。"""
    direction_map = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    shape_type = direction_map.get(direction, MSO_SHAPE.RIGHT_ARROW)
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _apply_fill_border(shape, fill_color, border_color, border_width)
    return shape


def add_star(
    slide,
    left, top, width, height,
    fill_color: RGBColor | str,
    points: int = 5,
    border_color: RGBColor | str | None = None,
    border_width: float = 0,
) -> object:
    """添加星形装饰（4/5/6/8/10/12/16/24/32 角星）。"""
    star_map = {
        4: MSO_SHAPE.STAR_4_POINT,
        5: MSO_SHAPE.STAR_5_POINT,
        6: MSO_SHAPE.STAR_6_POINT,
        7: MSO_SHAPE.STAR_7_POINT,
        8: MSO_SHAPE.STAR_8_POINT,
        10: MSO_SHAPE.STAR_10_POINT,
        12: MSO_SHAPE.STAR_12_POINT,
        16: MSO_SHAPE.STAR_16_POINT,
        24: MSO_SHAPE.STAR_24_POINT,
        32: MSO_SHAPE.STAR_32_POINT,
    }
    shape_type = star_map.get(points, MSO_SHAPE.STAR_5_POINT)
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _apply_fill_border(shape, fill_color, border_color, border_width)
    return shape


def add_chevron(
    slide,
    left, top, width, height,
    fill_color: RGBColor | str,
    direction: str = "right",
    border_color: RGBColor | str | None = None,
    border_width: float = 0,
) -> object:
    """添加V形箭头（流程/时间线用）。"""
    direction_map = {
        "right": MSO_SHAPE.CHEVRON,
        "left": MSO_SHAPE.LEFT_ARROW,  # 无原生左V形，用左箭头近似
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    shape_type = direction_map.get(direction, MSO_SHAPE.CHEVRON)
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _apply_fill_border(shape, fill_color, border_color, border_width)
    return shape


def add_pill_shape(
    slide,
    left, top, width, height,
    fill_color: RGBColor | str,
    border_color: RGBColor | str | None = None,
    border_width: float = 0,
) -> object:
    """添加胶囊形（圆角极值矩形，用于标签/大数字背景）。"""
    return add_rect(slide, left, top, width, height, fill_color, border_color, border_width, corner_radius=20)


def _apply_fill_border(shape, fill_color, border_color, border_width):
    """通用：为形状设置填充和边框。"""
    shape.fill.solid()
    if isinstance(fill_color, str):
        fill_color = _hex_to_rgb(fill_color)
    shape.fill.fore_color.rgb = fill_color
    if border_width > 0 and border_color:
        if isinstance(border_color, str):
            border_color = _hex_to_rgb(border_color)
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """将 '#RRGGBB' 或 'RRGGBB' 格式转为 RGBColor。"""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

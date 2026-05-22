"""卡片网格页面渲染器 — 3-4 个并列卡片，每个卡片含标题+描述。"""

from __future__ import annotations

from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_text_box, add_line, add_page_number,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
)


def render_grid_cards(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染卡片网格页面 — 将 body 中的每项渲染为独立卡片。"""
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    body = data.get("body", data.get("content", []))
    page_num = data.get("page_number", 0)
    total = data.get("_total_pages", 10)

    # 标题
    if title:
        add_text_box(
            slide,
            Inches(0.8), Inches(0.5),
            Inches(11.7), Inches(0.7),
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        add_line(
            slide,
            Inches(0.8), Inches(1.3),
            Inches(11.7), Inches(1.3),
            color=theme.accent,
            width=1.5,
        )

    # 卡片布局：每行最多 3 张卡片
    card_count = min(len(body), 6)
    cards_per_row = min(card_count, 3)
    rows = (card_count + cards_per_row - 1) // cards_per_row

    card_w = Inches(3.6)
    card_h = Inches(2.4)
    card_gap_x = Inches(0.4)
    card_gap_y = Inches(0.5)
    start_x = Inches(0.8)
    start_y = Inches(1.8)

    for i, item in enumerate(body[:card_count]):
        row = i // cards_per_row
        col = i % cards_per_row
        cx = start_x + col * (card_w + card_gap_x)
        cy = start_y + row * (card_h + card_gap_y)

        # 卡片背景
        add_rect(
            slide, cx, cy, card_w, card_h,
            fill_color=theme.light,
            border_color=theme.secondary,
            border_width=0.5,
            corner_radius=dna.corner_radius,
        )

        # 卡片顶部色条
        add_rect(
            slide, cx, cy, card_w, Inches(0.06),
            fill_color=theme.accent if col % 2 == 0 else theme.primary,
        )

        # 解析卡片内容（格式："标题 | 描述" 或 "标题: 描述" 或纯文本）
        card_text = str(item).strip()
        if " | " in card_text:
            card_title, card_desc = card_text.split(" | ", 1)
        elif ": " in card_text or "：" in card_text:
            sep = ": " if ": " in card_text else "："
            card_title, card_desc = card_text.split(sep, 1)
        else:
            card_title = card_text[:15]
            card_desc = card_text[15:] if len(card_text) > 15 else ""

        # 卡片标题
        add_text_box(
            slide,
            cx + Inches(0.3), cy + Inches(0.3),
            card_w - Inches(0.6), Inches(0.5),
            card_title,
            font_name=dna.title_font,
            font_size=18,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # 卡片描述
        if card_desc:
            add_text_box(
                slide,
                cx + Inches(0.3), cy + Inches(0.9),
                card_w - Inches(0.6), card_h - Inches(1.1),
                card_desc,
                font_name=dna.body_font,
                font_size=14,
                font_color=theme.text_body,
                alignment=PP_ALIGN.LEFT,
            )

    # 页码
    if page_num > 0 and page_num < total - 1:
        add_page_number(slide, page_num, total - 1, theme)

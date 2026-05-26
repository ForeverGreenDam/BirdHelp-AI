"""表格布局渲染器 — 结构化数据对比表格。

LLM 在 slide JSON 中提供 table_data 字段：
{
  "headers": ["列1标题", "列2标题", ...],
  "rows": [["单元格", "单元格", ...], ...],
  "highlight_col": 0 (可选，突出显示某列),
  "source": "数据来源（可选）"
}
"""

from __future__ import annotations

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_text_box, add_line, add_page_number,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
    _hex_to_rgb,
)
from loguru import logger


# 表格最大尺寸限制
_MAX_COLS = 8
_MAX_ROWS = 15


def render_table(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染表格页面：标题 + 结构化表格 + 来源标注。"""
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    table_data = data.get("table_data", {}) or {}
    page_num = data.get("page_number", 0)
    total = data.get("_total_pages", 10)

    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    highlight_col = table_data.get("highlight_col", 0)
    source = table_data.get("source", "")

    if not headers or not rows:
        _render_fallback(slide, data, theme, dna)
        return

    # 限制行列数
    headers = headers[:_MAX_COLS]
    rows = [[str(cell) for cell in row[:_MAX_COLS]] for row in rows[:_MAX_ROWS]]

    # ── 标题 ──
    if title:
        add_text_box(
            slide,
            Inches(0.8), Inches(0.4),
            Inches(11.7), Inches(0.7),
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        add_line(
            slide,
            Inches(0.8), Inches(1.2),
            Inches(11.7), Inches(1.2),
            color=theme.accent,
            width=1.5,
        )

    # ── 计算表格尺寸 ──
    num_cols = len(headers)
    num_rows = len(rows)

    table_left = Inches(0.8)
    table_top = Inches(1.5)
    table_width = Inches(11.7)

    # 列宽均分
    col_width = int(table_width / num_cols)
    # 行高自适应（标题 0.45 英寸，数据行 0.38 英寸），但不超过可用空间
    header_row_h = Inches(0.5)
    data_row_h = Inches(0.38)
    available_h = Inches(6.5) - table_top  # 到来源标注的距离
    max_table_h = min(
        header_row_h + data_row_h * num_rows + Inches(0.05),
        available_h,
    )
    # 实际数据行高
    if header_row_h + data_row_h * num_rows > max_table_h:
        data_row_h = int((max_table_h - header_row_h) / max(num_rows, 1))

    # ── 创建表格 ──
    table_shape = slide.shapes.add_table(
        num_rows + 1, num_cols,  # +1 为表头行
        table_left, table_top,
        table_width, header_row_h + data_row_h * num_rows,
    )
    pptx_table = table_shape.table

    # 设置列宽
    for ci in range(num_cols):
        pptx_table.columns[ci].width = col_width

    # ── 表头样式 ──
    header_fill = theme.table_header_fill or theme.primary
    for ci, header_text in enumerate(headers):
        cell = pptx_table.cell(0, ci)
        _set_cell_style(
            cell, str(header_text),
            fill_color=header_fill,
            font_color=RGBColor(255, 255, 255),
            font_name=dna.title_font,
            font_size=12,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

    # ── 数据行样式 ──
    alt_fill = theme.table_alt_fill or theme.light
    for ri, row in enumerate(rows):
        row_fill = theme.background if ri % 2 == 0 else alt_fill
        for ci, cell_text in enumerate(row):
            cell = pptx_table.cell(ri + 1, ci)
            is_first_col = ci == highlight_col
            _set_cell_style(
                cell, str(cell_text),
                fill_color=row_fill,
                font_color=theme.text_body,
                font_name=dna.body_font,
                font_size=11,
                bold=is_first_col,
                alignment=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
            )

    # ── 表格边框 ──
    border_color = theme.table_border or theme.secondary
    _set_table_borders(pptx_table, border_color)

    # ── 来源标注 ──
    source_top = table_top + header_row_h + data_row_h * num_rows + Inches(0.25)
    if source and source_top < Inches(6.8):
        add_text_box(
            slide,
            Inches(0.8), source_top,
            Inches(11.7), Inches(0.3),
            source,
            font_name=dna.body_font,
            font_size=9,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.LEFT,
        )

    # ── 页码 ──
    if page_num > 0 and page_num < total - 1:
        add_page_number(slide, page_num, total - 1, theme)


def _set_cell_style(
    cell,
    text: str,
    fill_color,
    font_color,
    font_name: str,
    font_size: int,
    bold: bool,
    alignment=PP_ALIGN.CENTER,  # noqa: ANN001
):
    """设置表格单元格的填充、字体和对齐。"""
    # 填充
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color

    # 文字
    tf = cell.text_frame
    tf.word_wrap = True
    # 垂直居中
    try:
        tf._txBody.bodyPr.set('anchor', 'ctr')
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color

    # 边距
    cell.margin_left = Pt(6)
    cell.margin_right = Pt(6)
    cell.margin_top = Pt(3)
    cell.margin_bottom = Pt(3)


def _set_table_borders(table, border_color):
    """为整个表格设置细边框。"""
    try:
        for ri in range(len(table.rows)):
            for ci in range(len(table.columns)):
                cell = table.cell(ri, ci)
                tcPr = cell._tc.get_or_add_tcPr()
                for border_name in ('lnL', 'lnR', 'lnT', 'lnB'):
                    ln_elem = tcPr.makeelement(qn(f'a:{border_name}'), {
                        'w': '6350',         # 0.5pt in EMU
                        'cap': 'flat',
                        'cmpd': 'sng',
                    })
                    # 边框颜色
                    solidFill = ln_elem.makeelement(qn('a:solidFill'), {})
                    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': _rgb_to_hex(border_color)})
                    solidFill.append(srgbClr)
                    ln_elem.append(solidFill)
                    # 移除旧边框再添加
                    existing = tcPr.find(qn(f'a:{border_name}'))
                    if existing is not None:
                        tcPr.remove(existing)
                    tcPr.append(ln_elem)
    except Exception:
        pass  # 边框失败不阻塞


def _rgb_to_hex(color) -> str:
    """RGBColor -> 'RRGGBB' hex 字符串。"""
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _render_fallback(slide, data, theme, dna):
    """表格数据缺失时回退为纯文字渲染。"""
    logger.warning("Table slide has no table_data, falling back to text_only")
    from generator.ppt.layouts.text_only import render_text_only
    render_text_only(slide, data, theme, dna, [])

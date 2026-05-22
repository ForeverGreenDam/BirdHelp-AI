"""章节过渡页渲染器 — 用于 PPT 中开启新话题的过渡页。"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_accent_bar, add_text_box, add_line, add_image,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
)


def render_section_header(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染章节分隔页。"""
    clear_placeholders(slide)
    variant = dna.section_variant
    image_path = images[0] if images else None

    if variant == 0:
        _render_section_left_bar(slide, data, theme, dna, image_path)
    elif variant == 1:
        _render_section_dark(slide, data, theme, dna, image_path)
    else:
        _render_section_numbered(slide, data, theme, dna, image_path)


def _render_section_left_bar(slide, data: dict, theme: ColorTheme, dna: DesignDNA,
                             image_path: str | None = None) -> None:
    """变体 0: 左侧强调条 + 章节标题。"""
    set_slide_bg(slide, theme.background)

    # 左侧强调条
    add_accent_bar(
        slide,
        Inches(0.8), Inches(2.0),
        Inches(0.1), Inches(3.5),
        color=theme.accent,
    )

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")

    text_width = Inches(7.0) if image_path else Inches(10.0)

    if title:
        add_text_box(
            slide,
            Inches(1.5), Inches(2.4),
            text_width, Inches(1.5),
            title,
            font_name=dna.title_font,
            font_size=38,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

    if subtitle:
        add_text_box(
            slide,
            Inches(1.5), Inches(4.2),
            text_width, Inches(0.8),
            subtitle,
            font_name=dna.body_font,
            font_size=18,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.LEFT,
        )

    # 底部分隔线
    add_line(
        slide,
        Inches(1.5), Inches(6.5),
        Inches(10.0), Inches(6.5),
        color=theme.light,
        width=1.0,
    )

    if image_path:
        try:
            add_image(slide, Inches(8.8), Inches(2.0), Inches(3.8), Inches(4.2), image_path)
        except Exception as exc:
            from loguru import logger
            logger.warning(f"section_header add_image failed: {exc}")


def _render_section_dark(slide, data: dict, theme: ColorTheme, dna: DesignDNA,
                         image_path: str | None = None) -> None:
    """变体 1: 深色全幅背景。"""
    set_slide_bg(slide, theme.dark)

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")

    title_top = Inches(2.0) if image_path else Inches(2.8)

    if image_path:
        try:
            add_image(slide, Inches(3.0), Inches(5.4), Inches(7.3), Inches(1.8), image_path)
        except Exception as exc:
            from loguru import logger
            logger.warning(f"section_dark add_image failed: {exc}")

    if title:
        add_text_box(
            slide,
            Inches(1.5), title_top,
            Inches(10.3), Inches(1.5),
            title,
            font_name=dna.title_font,
            font_size=40,
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

    # 标题下方装饰线
    add_rect(
        slide,
        Inches(5.8), title_top + Inches(1.8),
        Inches(1.7), Inches(0.04),
        fill_color=theme.accent,
    )

    if subtitle:
        add_text_box(
            slide,
            Inches(2.0), title_top + Inches(2.2),
            Inches(9.3), Inches(0.6),
            subtitle,
            font_name=dna.body_font,
            font_size=16,
            font_color=RGBColor(0xCC, 0xCC, 0xCC),
            alignment=PP_ALIGN.CENTER,
        )


def _render_section_numbered(slide, data: dict, theme: ColorTheme, dna: DesignDNA,
                             image_path: str | None = None) -> None:
    """变体 2: 章节编号 + 标题。"""
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    page_num = data.get("page_number", 1)

    # 章节序号水印
    add_text_box(
        slide,
        Inches(1.0), Inches(1.0),
        Inches(3.0), Inches(2.0),
        f"0{page_num // 3 + 1}" if page_num > 0 else "01",
        font_name=dna.title_font,
        font_size=72,
        font_color=theme.light,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )

    text_width = Inches(7.0) if image_path else Inches(10.0)

    if title:
        add_text_box(
            slide,
            Inches(1.5), Inches(3.2),
            text_width, Inches(1.3),
            title,
            font_name=dna.title_font,
            font_size=36,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

    if subtitle:
        add_text_box(
            slide,
            Inches(1.5), Inches(4.8),
            text_width, Inches(0.6),
            subtitle,
            font_name=dna.body_font,
            font_size=16,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.LEFT,
        )

    if image_path:
        try:
            add_image(slide, Inches(8.8), Inches(1.5), Inches(3.8), Inches(5.2), image_path)
        except Exception as exc:
            from loguru import logger
            logger.warning(f"section_numbered add_image failed: {exc}")

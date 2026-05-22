"""纯文字页面渲染器 — 标题 + 要点列表 + 装饰性形状。"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_accent_bar, add_text_box, add_line,
    add_multiline_text_box, add_page_number, add_image,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
)


def render_text_only(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染纯文字页面：标题 + 装饰 + 要点。有图片时右侧展示配图。"""
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    body = data.get("body", data.get("content", []))
    page_num = data.get("page_number", 0)
    total = data.get("_total_pages", 10)
    visual_plan = data.get("visual_plan", {})
    image_path = images[0] if images else None

    # 有图片时文字区收窄，右侧留给图片
    text_left = Inches(1.0)
    text_width = Inches(7.0) if image_path else Inches(11.3)

    # 顶部标题栏
    if title:
        add_text_box(
            slide,
            text_left, Inches(0.5),
            text_width, Inches(0.8),
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

    # 标题下分隔线
    decorations = visual_plan.get("decorations", [])
    has_line = any(d.get("type") == "line" and d.get("position") == "below_title" for d in decorations)
    if has_line or dna.show_decorations:
        add_line(
            slide,
            text_left, Inches(1.45),
            text_left + text_width, Inches(1.45),
            color=theme.accent,
            width=2.0,
        )

    if subtitle:
        add_text_box(
            slide,
            text_left, Inches(1.6),
            text_width, Inches(0.6),
            subtitle,
            font_name=dna.body_font,
            font_size=16,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.LEFT,
        )

    # 侧边强调条
    if dna.show_decorations:
        has_bar = any(d.get("type") == "accent_bar" for d in decorations)
        if has_bar or dna.decoration_level == "rich":
            add_accent_bar(
                slide,
                Inches(0.25), Inches(0.5),
                Inches(0.06), Inches(SLIDE_H.inches - 1.0),
                color=theme.accent,
            )

    # 正文要点
    if body:
        body_start_y = Inches(2.3) if subtitle else Inches(1.8)
        body_h = Inches(4.5)
        lines = _build_body_lines(body, theme, dna)
        add_multiline_text_box(
            slide,
            text_left + Inches(0.3), body_start_y,
            text_width - Inches(0.3), body_h,
            lines,
            theme,
        )

    # 右侧配图
    if image_path:
        try:
            add_image(
                slide,
                Inches(8.8), Inches(1.5),
                Inches(3.8), Inches(5.0),
                image_path,
            )
        except Exception as exc:
            from loguru import logger
            logger.warning(f"text_only add_image failed: {exc}")

    # 页码
    if page_num > 0 and page_num < total - 1:
        add_page_number(slide, page_num, total - 1, theme)


def _build_body_lines(body: list[str], theme: ColorTheme, dna: DesignDNA) -> list[dict]:
    """将要点列表转为多行配置。"""
    lines = []
    for item in body:
        if not item.strip():
            continue
        lines.append({
            "text": item,
            "size": dna.body_font_size,
            "bold": False,
            "color": theme.text_body,
            "bullet": True,
            "space_after": 8,
        })
    return lines

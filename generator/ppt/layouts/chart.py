"""图表布局渲染器 — 支持柱状图/折线图/饼图/面积图。

LLM 在 slide JSON 中提供 chart_data 字段：
{
  "chart_type": "bar" | "line" | "pie" | "area",
  "chart_title": "图表标题（可选，置空则用页面title）",
  "categories": ["类别1", "类别2", ...],
  "series": [{"name": "系列名", "data": [值1, 值2, ...]}, ...],
  "y_axis_label": "Y轴标签（可选）",
  "show_legend": true/false,
  "source": "数据来源（可选）"
}
"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_text_box, add_line, add_page_number,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
    _hex_to_rgb,
)
from loguru import logger

# 图表类型映射
_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    # 变体
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "line_stacked": XL_CHART_TYPE.LINE_MARKERS_STACKED,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


def render_chart(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染图表页面：标题 + 数据图表 + 来源标注。"""
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    chart_data = data.get("chart_data", {}) or {}
    page_num = data.get("page_number", 0)
    total = data.get("_total_pages", 10)

    chart_type_str = chart_data.get("chart_type", "bar")
    chart_type = _CHART_TYPE_MAP.get(chart_type_str)
    if chart_type is None:
        logger.warning(f"Unknown chart type '{chart_type_str}', falling back to bar")
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

    categories = chart_data.get("categories", [])
    series_list = chart_data.get("series", [])
    y_axis_label = chart_data.get("y_axis_label", "")
    show_legend = chart_data.get("show_legend", True)
    source = chart_data.get("source", "")

    if not categories or not series_list:
        # 无数据时渲染为纯文字页
        _render_fallback(slide, data, theme, dna)
        return

    # ── 标题 ──
    title_top = Inches(0.4)
    title_h = Inches(0.7)
    if title:
        add_text_box(
            slide,
            Inches(0.8), title_top,
            Inches(11.7), title_h,
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        add_line(
            slide,
            Inches(0.8), Inches(1.2),
            Inches(11.7), Inches(1.2),
            color=theme.accent,
            width=1.5,
        )

    # ── 构建图表数据 ──
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = categories
    for s in series_list:
        chart_data_obj.add_series(s.get("name", ""), s.get("data", []))

    # ── 图表位置与尺寸 ──
    chart_left = Inches(0.8)
    chart_top = Inches(1.5)
    chart_w = Inches(11.7)
    chart_h = Inches(4.8) if source else Inches(5.2)

    chart_frame = slide.shapes.add_chart(
        chart_type, chart_left, chart_top, chart_w, chart_h,
        chart_data_obj,
    )
    chart = chart_frame.chart

    # ── 应用主题颜色 ──
    chart_colors = theme.get_chart_colors()
    for i, series in enumerate(chart.series):
        color = chart_colors[i % len(chart_colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    # ── Y轴标签 ──
    if y_axis_label and hasattr(chart, 'value_axis'):
        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.paragraphs[0].text = y_axis_label
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        except Exception:
            pass

    # ── 图例 ──
    if show_legend and chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)

    # ── 数据标签 ──
    try:
        plot = chart.plots[0] if chart.plots else None
        for series in chart.series:
            # 直接启用数据标签（不依赖 has_data_labels，默认为 False）
            dl = series.data_labels
            dl.show_value = True
            dl.font.size = Pt(9)
            if chart_type_str in ("pie", "doughnut"):
                dl.show_percentage = True
                dl.show_category_name = True
                dl.show_value = False  # 饼图用百分比代替值
            elif chart_type_str in ("bar", "bar_stacked"):
                dl.show_value = True
                dl.font.size = Pt(8)
            elif chart_type_str == "line":
                dl.show_value = True
                dl.font.size = Pt(8)
    except Exception as exc:
        logger.debug(f"Data labels setup skipped: {exc}")

    # ── 来源标注 ──
    if source:
        add_text_box(
            slide,
            Inches(0.8), Inches(6.5),
            Inches(11.7), Inches(0.4),
            source,
            font_name=dna.body_font,
            font_size=9,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.LEFT,
        )

    # ── 页码 ──
    if page_num > 0 and page_num < total - 1:
        add_page_number(slide, page_num, total - 1, theme)


def _render_fallback(slide, data, theme, dna):
    """图表数据缺失时回退为纯文字渲染。"""
    from generator.ppt.layouts.text_only import render_text_only
    logger.warning(f"Chart slide has no chart_data, falling back to text_only")
    render_text_only(slide, data, theme, dna, [])

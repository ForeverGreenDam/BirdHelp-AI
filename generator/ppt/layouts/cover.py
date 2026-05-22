"""封面页渲染器 — 三种变体：渐变+强调条 / 色块分区 / 居中极简。"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_accent_bar, add_text_box, add_image,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
)


def render_cover(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染封面页，根据 dna.cover_variant 选择变体。"""
    clear_placeholders(slide)
    variant = dna.cover_variant
    image_path = images[0] if images else None

    if variant == 0:
        _render_cover_gradient(slide, data, theme, dna, image_path)
    elif variant == 1:
        _render_cover_split(slide, data, theme, dna, image_path)
    else:
        _render_cover_minimal(slide, data, theme, dna, image_path)


def _render_cover_gradient(slide, data: dict, theme: ColorTheme, dna: DesignDNA,
                          image_path: str | None = None) -> None:
    """变体 0: 深色背景 + 左侧强调条 + 标题左对齐。"""
    # 深色背景
    set_slide_bg(slide, theme.dark)

    # 左侧强调条
    add_accent_bar(
        slide,
        Inches(0), Inches(0),
        Inches(0.15), SLIDE_H,
        color=theme.accent,
    )

    # 顶部装饰圆
    from generator.ppt.shapes import add_circle
    add_circle(slide, Inches(11.5), Inches(1.0), Inches(2.0), theme.primary)

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    body = data.get("body", data.get("content", []))

    # 有配图时文字区收窄
    text_width = Inches(7.0) if image_path else Inches(10.5)

    if title:
        add_text_box(
            slide,
            Inches(1.5), Inches(2.2),
            text_width, Inches(1.8),
            title,
            font_name=dna.title_font,
            font_size=44,
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

    if subtitle:
        add_text_box(
            slide,
            Inches(1.5), Inches(4.2),
            text_width, Inches(0.8),
            subtitle,
            font_name=dna.body_font,
            font_size=22,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.LEFT,
        )

    # 底部附加信息
    if body:
        body_text = "  |  ".join(body) if isinstance(body, list) else body
        add_text_box(
            slide,
            Inches(1.5), Inches(5.5),
            text_width, Inches(0.6),
            body_text,
            font_name=dna.body_font,
            font_size=14,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.LEFT,
        )

    # 右侧配图
    if image_path:
        try:
            add_image(slide, Inches(8.8), Inches(1.5), Inches(3.8), Inches(4.8), image_path)
        except Exception as exc:
            from loguru import logger
            logger.warning(f"cover add_image failed: {exc}")


def _render_cover_split(slide, data: dict, theme: ColorTheme, dna: DesignDNA,
                       image_path: str | None = None) -> None:
    """变体 1: 左侧色块分区 + 右半白色。"""
    set_slide_bg(slide, theme.background)

    # 左侧大色块
    add_rect(
        slide,
        Inches(0), Inches(0),
        Inches(5.5), SLIDE_H,
        fill_color=theme.primary,
    )

    # 色块内装饰线
    add_rect(
        slide,
        Inches(3.5), Inches(2.5),
        Inches(0.08), Inches(2.5),
        fill_color=theme.accent,
    )

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    body = data.get("body", data.get("content", []))

    if title:
        add_text_box(
            slide,
            Inches(0.8), Inches(2.2),
            Inches(4.0), Inches(2.0),
            title,
            font_name=dna.title_font,
            font_size=40,
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

    # 右半区：有图片时上方放图，下方放文字
    right_text_top = Inches(2.8)

    if image_path:
        try:
            add_image(slide, Inches(6.5), Inches(1.0), Inches(5.5), Inches(3.0), image_path)
            right_text_top = Inches(4.3)
        except Exception as exc:
            from loguru import logger
            logger.warning(f"cover_split add_image failed: {exc}")

    if subtitle:
        add_text_box(
            slide,
            Inches(6.5), right_text_top,
            Inches(5.5), Inches(1.2),
            subtitle,
            font_name=dna.body_font,
            font_size=22,
            font_color=theme.text_primary,
            alignment=PP_ALIGN.LEFT,
        )

    if body:
        body_text = "\n".join(body) if isinstance(body, list) else body
        add_text_box(
            slide,
            Inches(6.5), right_text_top + Inches(1.7),
            Inches(5.5), Inches(1.5),
            body_text,
            font_name=dna.body_font,
            font_size=16,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.LEFT,
        )


def _render_cover_minimal(slide, data: dict, theme: ColorTheme, dna: DesignDNA,
                         image_path: str | None = None) -> None:
    """变体 2: 极简 — 白色背景 + 居中大标题 + 一根强调线。"""
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    body = data.get("body", data.get("content", []))

    # 有配图时整体上移为图片留空间
    title_top = Inches(1.2) if image_path else Inches(2.0)

    if image_path:
        try:
            add_image(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(1.6), image_path)
        except Exception as exc:
            from loguru import logger
            logger.warning(f"cover_minimal add_image failed: {exc}")

    if title:
        add_text_box(
            slide,
            Inches(1.5), title_top,
            Inches(10.3), Inches(2.0),
            title,
            font_name=dna.title_font,
            font_size=48,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

    # 标题下方细线
    add_rect(
        slide,
        Inches(5.5), title_top + Inches(2.3),
        Inches(2.3), Inches(0.04),
        fill_color=theme.accent,
    )

    if subtitle:
        add_text_box(
            slide,
            Inches(2.0), title_top + Inches(2.8),
            Inches(9.3), Inches(0.8),
            subtitle,
            font_name=dna.body_font,
            font_size=20,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.CENTER,
        )

    if body:
        body_text = "  |  ".join(body) if isinstance(body, list) else body
        add_text_box(
            slide,
            Inches(2.0), title_top + Inches(3.8),
            Inches(9.3), Inches(0.6),
            body_text,
            font_name=dna.body_font,
            font_size=14,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.CENTER,
        )

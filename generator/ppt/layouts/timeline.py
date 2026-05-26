"""时间线布局渲染器 — 发展历程/里程碑/项目阶段展示。

body 格式：["时间 | 标题 | 描述", ...]
每条按时间顺序自上而下或自左向右排列。
"""

from __future__ import annotations

from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_circle, add_line, add_text_box, add_page_number, add_chevron,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
    _hex_to_rgb,
)


def render_timeline(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染时间线页面：标题 + 多个里程碑节点连成的水平时间线。

    每个节点格式 "时间 | 标题 | 描述"，最多 6 个节点。
    """
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    body = data.get("body", data.get("content", []))
    page_num = data.get("page_number", 0)
    total = data.get("_total_pages", 10)

    nodes = _parse_nodes(body)

    # ── 标题 ──
    if title:
        add_text_box(
            slide,
            Inches(0.8), Inches(0.4),
            Inches(11.7), Inches(0.7),
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        add_line(
            slide,
            Inches(0.8), Inches(1.2),
            Inches(11.7), Inches(1.2),
            color=theme.accent,
            width=1.5,
        )

    if not nodes:
        add_text_box(
            slide, Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.0),
            "时间线数据为空", font_name=dna.body_font, font_size=18,
            font_color=theme.text_secondary, alignment=PP_ALIGN.CENTER,
        )
        return

    num_nodes = min(len(nodes), 5)

    # ── 水平时间线 ──
    line_y = Inches(3.5)
    margin = Inches(1.2)
    usable_w = SLIDE_W - margin * 2
    # 动态间距：节点越多间距越小
    if num_nodes <= 2:
        node_spacing = Emu(int(usable_w / 2))
        start_x = margin + Inches(1.5)
    elif num_nodes == 3:
        node_spacing = Emu(int(usable_w / 2.5))
        start_x = margin + Inches(0.8)
    else:
        node_spacing = Emu(int(usable_w / (num_nodes - 1)))
        start_x = margin

    # 根据节点数动态调整文字框宽度，防止重叠
    text_box_w = Inches(min(2.2, node_spacing.inches * 0.85))
    time_box_w = Inches(min(1.6, node_spacing.inches * 0.7))

    # 绘制主线
    add_line(
        slide,
        margin, line_y,
        SLIDE_W - margin, line_y,
        color=theme.secondary,
        width=2.5,
    )

    # ── 节点（上下交错排列避免拥挤） ──
    for i, (time_str, node_title, node_desc) in enumerate(nodes[:num_nodes]):
        nx = start_x + i * node_spacing
        is_top = i % 2 == 0

        # 节点圆
        circle_radius = Inches(0.14)
        add_circle(slide, nx, line_y, circle_radius, theme.accent)

        # 连接竖线（增加长度给文字更多空间）
        connector_len = Inches(0.7)
        if is_top:
            cy = line_y - connector_len
            add_line(slide, nx, line_y - circle_radius, nx, cy, theme.secondary, 1.0)
        else:
            cy = line_y + connector_len
            add_line(slide, nx, line_y + circle_radius, nx, cy, theme.secondary, 1.0)

        # 动态字号
        time_fs = 13 if num_nodes <= 4 else 11
        title_fs = 12 if num_nodes <= 4 else 10

        # 每个文字框居中对齐于节点
        half_tw = int(time_box_w / 2)
        half_bw = int(text_box_w / 2)

        if is_top:
            # 时间标签
            add_text_box(slide, nx - half_tw, cy - Inches(0.5),
                         time_box_w, Inches(0.3),
                         time_str, font_name=dna.body_font, font_size=time_fs,
                         font_color=theme.accent, bold=True, alignment=PP_ALIGN.CENTER)
            # 标题
            add_text_box(slide, nx - half_bw, cy - Inches(0.9),
                         text_box_w, Inches(0.35),
                         node_title, font_name=dna.title_font, font_size=title_fs,
                         font_color=theme.text_primary, bold=True, alignment=PP_ALIGN.CENTER)
            # 描述
            if node_desc and num_nodes <= 4:
                add_text_box(slide, nx - half_bw, cy - Inches(1.3),
                             text_box_w, Inches(0.5),
                             node_desc, font_name=dna.body_font, font_size=9,
                             font_color=theme.text_secondary, alignment=PP_ALIGN.CENTER)
        else:
            add_text_box(slide, nx - half_tw, cy + Inches(0.15),
                         time_box_w, Inches(0.3),
                         time_str, font_name=dna.body_font, font_size=time_fs,
                         font_color=theme.accent, bold=True, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, nx - half_bw, cy + Inches(0.5),
                         text_box_w, Inches(0.35),
                         node_title, font_name=dna.title_font, font_size=title_fs,
                         font_color=theme.text_primary, bold=True, alignment=PP_ALIGN.CENTER)
            if node_desc and num_nodes <= 4:
                add_text_box(slide, nx - half_bw, cy + Inches(0.85),
                             text_box_w, Inches(0.5),
                             node_desc, font_name=dna.body_font, font_size=9,
                             font_color=theme.text_secondary, alignment=PP_ALIGN.CENTER)

    # ── 页码 ──
    if page_num > 0 and page_num < total - 1:
        add_page_number(slide, page_num, total - 1, theme)


def _parse_nodes(body: list) -> list[tuple[str, str, str]]:
    """解析 '时间 | 标题 | 描述' 格式为节点列表。"""
    if isinstance(body, str):
        body = [body]
    nodes = []
    for item in body:
        text = str(item).strip()
        parts = [p.strip() for p in text.split("|")]
        if len(parts) >= 3:
            nodes.append((parts[0], parts[1], parts[2]))
        elif len(parts) == 2:
            nodes.append((parts[0], parts[1], ""))
        elif len(parts) == 1:
            nodes.append(("", parts[0], ""))
    return nodes

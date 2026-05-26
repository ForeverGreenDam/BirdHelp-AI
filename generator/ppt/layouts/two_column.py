"""双栏布局页面渲染器 — 左右两栏对比或并列展示。"""

from __future__ import annotations

from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_accent_bar, add_text_box, add_line,
    add_multiline_text_box, add_page_number, add_image,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
)


def render_two_column(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染双栏对比/并列页面。"""
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    body = data.get("body", data.get("content", []))
    page_num = data.get("page_number", 0)
    total = data.get("_total_pages", 10)

    # 标题
    if title:
        add_text_box(
            slide,
            Inches(0.8), Inches(0.5),
            Inches(11.7), Inches(0.8),
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        add_line(
            slide,
            Inches(0.8), Inches(1.4),
            Inches(11.7), Inches(1.4),
            color=theme.accent,
            width=1.5,
        )

    # 将内容分为左右两栏
    # 当 body 恰为 2 个元素时视为 [左栏内容, 右栏内容]，每项按换行拆分为要点
    if len(body) == 2:
        left_raw, right_raw = body[0], body[1]
        left_items = [l.strip() for l in str(left_raw).split("\n") if l.strip()]
        right_items = [r.strip() for r in str(right_raw).split("\n") if r.strip()]
    else:
        mid = len(body) // 2 if len(body) > 1 else 1
        left_items = body[:mid]
        right_items = body[mid:]

    col_w = Inches(5.6)
    image_path = images[0] if images else None
    # 有图片时缩短列高，为底部图片留空间
    col_top = Inches(1.8)
    col_h = Inches(3.5) if image_path else Inches(4.8)

    # 左栏
    if left_items:
        # 左栏标题背板
        left_label = data.get("left_label", "")
        if left_label:
            add_rect(
                slide,
                Inches(0.8), col_top,
                col_w, Inches(0.5),
                fill_color=theme.primary,
            )
            add_text_box(
                slide,
                Inches(1.0), col_top + Inches(0.05),
                col_w - Inches(0.4), Inches(0.4),
                left_label,
                font_name=dna.body_font,
                font_size=15,
                font_color="FFFFFF",
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            left_body_top = col_top + Inches(0.7)
        else:
            left_body_top = col_top

        left_lines = _build_body_lines(left_items, theme, dna)
        add_multiline_text_box(
            slide,
            Inches(1.0), left_body_top,
            col_w - Inches(0.4), col_h - (left_body_top - col_top),
            left_lines,
            theme,
        )

    # 中间分割线
    add_line(
        slide,
        Inches(6.67), Inches(1.8),
        Inches(6.67), Inches(6.8),
        color=theme.light,
        width=1.5,
    )

    # 右栏
    if right_items:
        right_label = data.get("right_label", "")
        if right_label:
            add_rect(
                slide,
                Inches(7.2), col_top,
                col_w, Inches(0.5),
                fill_color=theme.accent,
            )
            add_text_box(
                slide,
                Inches(7.4), col_top + Inches(0.05),
                col_w - Inches(0.4), Inches(0.4),
                right_label,
                font_name=dna.body_font,
                font_size=15,
                font_color="FFFFFF",
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            right_body_top = col_top + Inches(0.7)
        else:
            right_body_top = col_top

        right_lines = _build_body_lines(right_items, theme, dna)
        add_multiline_text_box(
            slide,
            Inches(7.4), right_body_top,
            col_w - Inches(0.4), col_h - (right_body_top - col_top),
            right_lines,
            theme,
        )

    # 底部配图
    if image_path:
        try:
            img_top = col_top + col_h + Inches(0.3)
            img_h = SLIDE_H - img_top - Inches(0.5)
            add_image(slide, Inches(1.5), img_top, Inches(10.3), img_h, image_path)
        except Exception:
            pass

    # 页码
    if page_num > 0 and page_num < total - 1:
        add_page_number(slide, page_num, total - 1, theme)


def _build_body_lines(body: list[str], theme: ColorTheme, dna: DesignDNA) -> list[dict]:
    if isinstance(body, str):
        body = [body]
    lines = []
    for item in body:
        item = str(item).strip()
        if not item:
            continue
        lines.append({
            "text": item,
            "size": dna.body_font_size,
            "bold": False,
            "color": theme.text_body,
            "bullet": True,
            "space_after": 8,
        })
    return lines

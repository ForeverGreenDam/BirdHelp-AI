"""大数字展示布局 — 核心KPI/指标突出显示。

body 格式：["数字 | 标签", ...]
每条拆分为大数字行和标签行，最多展示 6 个指标。
"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_text_box, add_line, add_page_number, add_pill_shape,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
    _hex_to_rgb,
)


def render_big_number(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染大数字页面：标题 + 2-4 个大数字指标。"""
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    body = data.get("body", data.get("content", []))
    page_num = data.get("page_number", 0)
    total = data.get("_total_pages", 10)

    # ── 解析指标 ──
    metrics = _parse_metrics(body)
    if not metrics:
        metrics = [("—", "无数据")]

    num_metrics = min(len(metrics), 6)
    metrics = metrics[:num_metrics]

    # ── 标题 ──
    title_y = Inches(0.4)
    if title:
        add_text_box(
            slide,
            Inches(0.8), title_y,
            Inches(11.7), Inches(0.7),
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        add_line(
            slide,
            Inches(0.8), Inches(1.2),
            Inches(11.7), Inches(1.2),
            color=theme.accent,
            width=1.5,
        )

    # ── 指标卡片布局 ──
    # 根据数量自适应布局：1居中，2并排，3一排，4的2x2网格，5/6的2行网格
    if num_metrics == 1:
        _render_single(slide, metrics[0], theme, dna)
    elif num_metrics == 2:
        _render_two(slide, metrics, theme, dna)
    elif num_metrics == 3:
        _render_three(slide, metrics, theme, dna)
    elif num_metrics == 4:
        _render_four(slide, metrics, theme, dna)
    elif num_metrics == 5:
        _render_five(slide, metrics, theme, dna)
    else:
        _render_six(slide, metrics, theme, dna)

    # ── 页码 ──
    if page_num > 0 and page_num < total - 1:
        add_page_number(slide, page_num, total - 1, theme)


def _parse_metrics(body: list) -> list[tuple[str, str]]:
    """解析 body 中的 '数字 | 标签' 格式。兼容 LLM 修复时输出的字符串格式。"""
    if isinstance(body, str):
        body = [body]
    metrics = []
    for item in body:
        text = str(item).strip()
        if " | " in text:
            num, label = text.split(" | ", 1)
        elif "|" in text:
            num, label = text.split("|", 1)
        elif "：" in text or ": " in text:
            sep = "：" if "：" in text else ": "
            num, label = text.split(sep, 1)
        else:
            num = text
            label = ""
        metrics.append((num.strip(), label.strip()))
    return metrics


def _render_single(slide, metric, theme, dna):
    """单个大数字居中展示。"""
    num, label = metric
    cx = Inches(6.67)
    cy = Inches(3.5)

    # 背景 pill
    add_pill_shape(
        slide,
        Inches(3.0), Inches(2.2),
        Inches(7.3), Inches(2.8),
        theme.light,
    )

    # 大数字
    add_text_box(
        slide,
        Inches(3.5), Inches(2.5),
        Inches(6.3), Inches(1.3),
        num,
        font_name=dna.title_font,
        font_size=56,
        font_color=theme.accent,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    # 标签
    if label:
        add_text_box(
            slide,
            Inches(3.5), Inches(3.9),
            Inches(6.3), Inches(0.6),
            label,
            font_name=dna.body_font,
            font_size=18,
            font_color=theme.text_body,
            alignment=PP_ALIGN.CENTER,
        )


def _render_two(slide, metrics, theme, dna):
    """两个大数字并排展示。"""
    card_w = Inches(5.0)
    card_h = Inches(2.5)
    start_x = Inches(1.67)
    cy = Inches(2.8)

    for i, (num, label) in enumerate(metrics):
        cx = start_x + i * (card_w + Inches(0.6))
        # 卡片背景
        add_rect(
            slide, cx, cy, card_w, card_h,
            theme.light, border_color=theme.secondary, border_width=0.5,
            corner_radius=dna.corner_radius,
        )
        # 顶部强调色条
        add_rect(slide, cx, cy, card_w, Inches(0.05), theme.accent)
        # 数字
        add_text_box(
            slide, cx + Inches(0.3), cy + Inches(0.4),
            card_w - Inches(0.6), Inches(1.1),
            num, font_name=dna.title_font, font_size=48,
            font_color=theme.accent, bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        # 标签
        if label:
            add_text_box(
                slide, cx + Inches(0.3), cy + Inches(1.6),
                card_w - Inches(0.6), Inches(0.5),
                label, font_name=dna.body_font, font_size=16,
                font_color=theme.text_body,
                alignment=PP_ALIGN.CENTER,
            )


def _render_three(slide, metrics, theme, dna):
    """三个大数字并排展示。"""
    card_w = Inches(3.4)
    card_h = Inches(2.5)
    start_x = Inches(1.1)
    cy = Inches(2.8)

    for i, (num, label) in enumerate(metrics):
        cx = start_x + i * (card_w + Inches(0.35))
        add_rect(
            slide, cx, cy, card_w, card_h,
            theme.light, border_color=theme.secondary, border_width=0.5,
            corner_radius=dna.corner_radius,
        )
        add_rect(slide, cx, cy, card_w, Inches(0.05), theme.accent)
        add_text_box(
            slide, cx + Inches(0.2), cy + Inches(0.4),
            card_w - Inches(0.4), Inches(1.1),
            num, font_name=dna.title_font, font_size=42,
            font_color=theme.accent, bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        if label:
            add_text_box(
                slide, cx + Inches(0.2), cy + Inches(1.6),
                card_w - Inches(0.4), Inches(0.5),
                label, font_name=dna.body_font, font_size=14,
                font_color=theme.text_body,
                alignment=PP_ALIGN.CENTER,
            )


def _render_four(slide, metrics, theme, dna):
    """四个大数字 2x2 网格展示。"""
    card_w = Inches(5.0)
    card_h = Inches(2.1)
    start_x = Inches(1.67)
    col_gap = Inches(0.6)
    row_gap = Inches(0.5)
    start_y = Inches(2.0)

    for i, (num, label) in enumerate(metrics):
        col = i % 2
        row = i // 2
        cx = start_x + col * (card_w + col_gap)
        cy = start_y + row * (card_h + row_gap)
        add_rect(
            slide, cx, cy, card_w, card_h,
            theme.light, border_color=theme.secondary, border_width=0.5,
            corner_radius=dna.corner_radius,
        )
        add_rect(slide, cx, cy, card_w, Inches(0.05), theme.accent)
        add_text_box(
            slide, cx + Inches(0.3), cy + Inches(0.25),
            card_w - Inches(0.6), Inches(0.9),
            num, font_name=dna.title_font, font_size=38,
            font_color=theme.accent, bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        if label:
            add_text_box(
                slide, cx + Inches(0.3), cy + Inches(1.3),
                card_w - Inches(0.6), Inches(0.5),
                label, font_name=dna.body_font, font_size=14,
                font_color=theme.text_body,
                alignment=PP_ALIGN.CENTER,
            )


def _render_five(slide, metrics, theme, dna):
    """五个大数字展示：第一行 3 个，第二行 2 个居中。"""
    card_w = Inches(3.6)
    card_h = Inches(2.0)
    gap_x = Inches(0.3)
    start_y = Inches(2.0)
    row_gap = Inches(0.4)

    # 第一行: 3 个卡片
    row0_width = 3 * card_w + 2 * gap_x
    start_x_0 = (SLIDE_W - row0_width) / 2
    for i in range(3):
        cx = start_x_0 + i * (card_w + gap_x)
        cy = start_y
        _draw_metric_card(slide, metrics[i], cx, cy, card_w, card_h,
                          theme, dna, font_size=34, label_size=13)

    # 第二行: 2 个卡片居中
    row1_width = 2 * card_w + gap_x
    start_x_1 = (SLIDE_W - row1_width) / 2
    for i in range(3, 5):
        cx = start_x_1 + (i - 3) * (card_w + gap_x)
        cy = start_y + card_h + row_gap
        _draw_metric_card(slide, metrics[i], cx, cy, card_w, card_h,
                          theme, dna, font_size=34, label_size=13)


def _render_six(slide, metrics, theme, dna):
    """六个大数字 2x3 网格展示。"""
    card_w = Inches(3.6)
    card_h = Inches(2.0)
    gap_x = Inches(0.35)
    start_y = Inches(1.9)
    row_gap = Inches(0.35)

    row_width = 3 * card_w + 2 * gap_x
    start_x = (SLIDE_W - row_width) / 2

    for i, _ in enumerate(metrics):
        col = i % 3
        row = i // 3
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + row_gap)
        _draw_metric_card(slide, metrics[i], cx, cy, card_w, card_h,
                          theme, dna, font_size=32, label_size=12)


def _draw_metric_card(slide, metric, cx, cy, card_w, card_h,
                      theme, dna, font_size=38, label_size=14):
    """绘制单个指标卡片。"""
    num, label = metric
    add_rect(
        slide, cx, cy, card_w, card_h,
        theme.light, border_color=theme.secondary, border_width=0.5,
        corner_radius=dna.corner_radius,
    )
    add_rect(slide, cx, cy, card_w, Inches(0.05), theme.accent)
    add_text_box(
        slide, cx + Inches(0.2), cy + Inches(0.2),
        card_w - Inches(0.4), Inches(0.8),
        num, font_name=dna.title_font, font_size=font_size,
        font_color=theme.accent, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    if label:
        add_text_box(
            slide, cx + Inches(0.2), cy + Inches(1.1),
            card_w - Inches(0.4), Inches(0.5),
            label, font_name=dna.body_font, font_size=label_size,
            font_color=theme.text_body,
            alignment=PP_ALIGN.CENTER,
        )

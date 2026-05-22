"""总结/致谢页渲染器 — PPT 最后一页。"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_accent_bar, add_text_box, add_line, add_circle, add_image,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
)


def render_summary(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染总结/致谢页。"""
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "感谢观看")
    body = data.get("body", data.get("content", []))
    subtitle = data.get("subtitle", "")
    image_path = images[0] if images else None

    # 顶部装饰条
    add_rect(
        slide,
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.08),
        fill_color=theme.accent,
    )

    # 有配图时整体上移
    title_top = Inches(1.2) if image_path else Inches(2.2)

    # 装饰圆
    if dna.show_decorations and dna.decoration_level != "minimal":
        add_circle(slide, Inches(6.67), Inches(6.2), Inches(1.8), theme.light)

    # 配图（居中横幅）
    if image_path:
        try:
            add_image(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(1.6), image_path)
        except Exception as exc:
            from loguru import logger
            logger.warning(f"summary add_image failed: {exc}")

    # 主标题（居中大字体）
    add_text_box(
        slide,
        Inches(1.5), title_top,
        Inches(10.3), Inches(1.5),
        title,
        font_name=dna.title_font,
        font_size=44,
        font_color=theme.text_primary,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # 标题下方装饰线
    add_line(
        slide,
        Inches(5.2), title_top + Inches(1.6),
        Inches(8.1), title_top + Inches(1.6),
        color=theme.accent,
        width=2.0,
    )

    if subtitle:
        add_text_box(
            slide,
            Inches(2.0), title_top + Inches(2.0),
            Inches(9.3), Inches(0.7),
            subtitle,
            font_name=dna.body_font,
            font_size=20,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.CENTER,
        )

    # 附加信息
    if body:
        body_text = "\n".join(body) if isinstance(body, list) else body
        add_text_box(
            slide,
            Inches(2.0), title_top + Inches(3.0),
            Inches(9.3), Inches(1.0),
            body_text,
            font_name=dna.body_font,
            font_size=15,
            font_color=theme.text_secondary,
            alignment=PP_ALIGN.CENTER,
        )

    # 底部装饰条
    add_rect(
        slide,
        Inches(0), Inches(7.42),
        SLIDE_W, Inches(0.08),
        fill_color=theme.accent,
    )

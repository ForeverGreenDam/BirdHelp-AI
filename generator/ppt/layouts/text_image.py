"""图文混排页面渲染器 — 左文右图 / 上图下文 两种模式。"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from generator.ppt.theme import ColorTheme
from generator.ppt.layout import DesignDNA
from generator.ppt.shapes import (
    add_rect, add_accent_bar, add_text_box, add_line,
    add_multiline_text_box, add_image, add_page_number,
    set_slide_bg, clear_placeholders, SLIDE_W, SLIDE_H,
)


def render_text_image(
    slide,
    data: dict,
    theme: ColorTheme,
    dna: DesignDNA,
    images: list[str],
) -> None:
    """渲染图文混排页面。根据 visual_plan.layout_hint 或 image_position 决定布局方向。"""
    clear_placeholders(slide)
    set_slide_bg(slide, theme.background)

    title = data.get("title", "")
    body = data.get("body", data.get("content", []))
    page_num = data.get("page_number", 0)
    total = data.get("_total_pages", 10)
    visual_plan = data.get("visual_plan", {})
    layout_hint = visual_plan.get("layout_hint", data.get("image_position", "right"))

    image_path = images[0] if images else None

    if layout_hint in ("top", "image_top_text_bottom"):
        _render_top_image(slide, title, body, theme, dna, image_path)
    elif layout_hint == "left" or (layout_hint == "right" and image_path):
        _render_side_by_side(slide, title, body, theme, dna, image_path, layout_hint)
    else:
        # 默认：左文右图
        _render_side_by_side(slide, title, body, theme, dna, image_path, "right")

    # 页码
    if page_num > 0 and page_num < total - 1:
        add_page_number(slide, page_num, total - 1, theme)


def _render_side_by_side(
    slide,
    title: str,
    body: list[str],
    theme: ColorTheme,
    dna: DesignDNA,
    image_path: str | None,
    image_side: str,
) -> None:
    """左右布局 — 一侧文字一侧图片。"""
    text_left = image_side == "right"
    text_x = Inches(0.8) if text_left else Inches(7.0)
    text_w = Inches(5.8)
    img_x = Inches(7.0) if text_left else Inches(0.8)
    img_w = Inches(5.8)

    # 标题
    if title:
        add_text_box(
            slide,
            text_x, Inches(0.6),
            text_w, Inches(0.8),
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        # 标题下短线
        add_line(
            slide,
            text_x, Inches(1.5),
            text_x + Inches(2.0), Inches(1.5),
            color=theme.accent,
            width=2.0,
        )

    # 正文
    if body:
        lines = _build_body_lines(body, theme, dna)
        add_multiline_text_box(
            slide,
            text_x, Inches(1.8),
            text_w, Inches(4.8),
            lines,
            theme,
        )

    # 图片
    if image_path:
        try:
            add_image(
                slide,
                img_x, Inches(1.2),
                img_w, Inches(5.2),
                image_path,
            )
        except Exception as exc:
            from loguru import logger
            logger.warning(f"add_image failed for slide, falling back to placeholder: {exc}")
            _draw_image_placeholder(slide, img_x, Inches(1.2), img_w, Inches(5.2), theme)
    else:
        _draw_image_placeholder(slide, img_x, Inches(1.2), img_w, Inches(5.2), theme)


def _render_top_image(
    slide,
    title: str,
    body: list[str],
    theme: ColorTheme,
    dna: DesignDNA,
    image_path: str | None,
) -> None:
    """上图下文布局。"""
    img_area_h = Inches(3.5)

    if image_path:
        try:
            add_image(slide, Inches(0.5), Inches(0.3), Inches(12.3), img_area_h, image_path)
        except Exception as exc:
            from loguru import logger
            logger.warning(f"add_image (top) failed for slide, falling back to placeholder: {exc}")
            _draw_image_placeholder(slide, Inches(0.5), Inches(0.3), Inches(12.3), img_area_h, theme)
    else:
        _draw_image_placeholder(slide, Inches(0.5), Inches(0.3), Inches(12.3), img_area_h, theme)

    # 标题
    if title:
        add_text_box(
            slide,
            Inches(0.8), Inches(4.0),
            Inches(11.7), Inches(0.7),
            title,
            font_name=dna.title_font,
            font_size=dna.title_font_size,
            font_color=theme.text_primary,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )

    # 正文
    if body:
        lines = _build_body_lines(body, theme, dna)
        add_multiline_text_box(
            slide,
            Inches(0.8), Inches(4.8),
            Inches(11.7), Inches(2.2),
            lines,
            theme,
        )


def _build_body_lines(body: list[str], theme: ColorTheme, dna: DesignDNA) -> list[dict]:
    """将要点列表转为多行配置。"""
    lines = []
    for item in body:
        if not item.strip():
            continue
        lines.append({
            "text": item,
            "size": dna.body_font_size,
            "bold": False,
            "color": theme.text_body,
            "bullet": True,
            "space_after": 8,
        })
    return lines


def _draw_image_placeholder(slide, left, top, width, height, theme: ColorTheme) -> None:
    """在图片缺失时绘制一个带提示文字的灰色占位区块。"""
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from generator.ppt.shapes import add_rect, add_text_box

    add_rect(
        slide, left, top, width, height,
        fill_color=theme.light,
        border_color=theme.secondary,
        border_width=1.0,
    )
    add_text_box(
        slide,
        left + Inches(0.5), top + height / 3,
        width - Inches(1.0), Inches(0.8),
        "[配图区域]",
        font_name=theme.body_font,
        font_size=14,
        font_color=theme.text_secondary,
        alignment=PP_ALIGN.CENTER,
    )

"""PPT 文件生成器 — 基于 python-pptx 将结构化大纲构建为 .pptx 文件。"""

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from loguru import logger

from generator.base import BaseGenerator

STYLE_THEMES = {
    "academic": {
        "title_color": RGBColor(0x1A, 0x3C, 0x6E),
        "subtitle_color": RGBColor(0x55, 0x6B, 0x8D),
        "body_color": RGBColor(0x2D, 0x2D, 0x2D),
        "accent_color": RGBColor(0x1A, 0x3C, 0x6E),
        "bg_light": RGBColor(0xF0, 0xF2, 0xF5),
        "title_font": "WenQuanYi Micro Hei",
        "body_font": "AR PL UMing CN",
    },
    "business": {
        "title_color": RGBColor(0x1B, 0x3A, 0x5C),
        "subtitle_color": RGBColor(0x5A, 0x72, 0x8C),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "accent_color": RGBColor(0x00, 0x6E, 0xB6),
        "bg_light": RGBColor(0xE8, 0xEC, 0xF0),
        "title_font": "WenQuanYi Micro Hei",
        "body_font": "WenQuanYi Micro Hei",
    },
    "creative": {
        "title_color": RGBColor(0xE0, 0x4A, 0x36),
        "subtitle_color": RGBColor(0xE8, 0x7A, 0x3C),
        "body_color": RGBColor(0x3C, 0x3C, 0x3C),
        "accent_color": RGBColor(0xE0, 0x4A, 0x36),
        "bg_light": RGBColor(0xFE, 0xF6, 0xF0),
        "title_font": "WenQuanYi Micro Hei",
        "body_font": "WenQuanYi Micro Hei",
    },
}

LAYOUT_MAP = {
    "title_slide": 0,
    "title_and_content": 1,
    "section_header": 2,
    "two_content": 3,
    "blank": 6,
}


class PptGenerator(BaseGenerator):
    """PPT 生成器，将 LLM 输出的结构化 JSON 大纲渲染为 .pptx 文件。"""

    output_extension = ".pptx"

    def generate(self, content: dict[str, Any], output_path: Path) -> Path:
        """根据结构化内容生成 PPT 文件，返回输出路径。"""
        parsed = self._parse_content(content)
        slides = parsed.get("slides", [])
        title = parsed.get("title", "演示文稿")
        style_name = parsed.get("style", "academic")
        theme = STYLE_THEMES.get(style_name, STYLE_THEMES["academic"])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, slide_data in enumerate(slides):
            self._add_slide(prs, slide_data, theme, i, len(slides))

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info(f"PPT generated: {output_path}, {len(slides)} slides, style={style_name}")
        return output_path

    def _add_slide(self, prs, slide_data: dict, theme: dict, index: int, total: int) -> None:
        """添加单页幻灯片。"""
        layout_name = slide_data.get("layout", "title_and_content")
        layout_idx = LAYOUT_MAP.get(layout_name, 1)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # 移除所有占位符，避免空占位符导致的"文件损坏"警告和内容被遮挡
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # 背景色
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = theme["bg_light"]

        # 标题 — 根据布局调整位置
        slide_title = slide_data.get("title", "")
        if slide_title:
            if layout_name == "title_slide":
                self._add_textbox(slide, slide_title,
                                  Inches(1.0), Inches(2.0),
                                  Inches(11.3), Inches(1.5), theme, "title")
            elif layout_name == "section_header":
                self._add_textbox(slide, slide_title,
                                  Inches(1.0), Inches(2.8),
                                  Inches(11.3), Inches(1.5), theme, "title")
            else:
                self._add_textbox(slide, slide_title,
                                  Inches(0.8), Inches(0.3),
                                  Inches(11.7), Inches(1.0), theme, "title")

        # 副标题
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            if layout_name == "title_slide":
                self._add_textbox(slide, subtitle,
                                  Inches(1.5), Inches(3.8),
                                  Inches(10.3), Inches(1.0), theme, "subtitle")
            else:
                self._add_textbox(slide, subtitle,
                                  Inches(1.0), Inches(1.4),
                                  Inches(11.3), Inches(0.8), theme, "subtitle")

        # 正文内容
        content_top = Inches(2.4) if subtitle else Inches(1.6)
        content_items = slide_data.get("content", [])
        if content_items:
            if layout_name == "two_content":
                mid = len(content_items) // 2
                self._fill_content(slide, content_items[:mid], theme, Inches(1.0), content_top,
                                   Inches(5.3), Inches(4.8))
                self._fill_content(slide, content_items[mid:], theme, Inches(7.0), content_top,
                                   Inches(5.3), Inches(4.8))
            else:
                self._fill_content(slide, content_items, theme, Inches(1.0), content_top,
                                   Inches(11.3), Inches(4.5))

        # 演讲备注
        notes_text = slide_data.get("notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

        # 页码
        if index > 0 and index < total - 1:
            self._add_page_number(slide, index, total, theme)

    def _fill_content(self, slide, items: list[str], theme: dict,
                      left, top, width, height) -> None:
        """在指定区域填充要点列表。left/top/width/height 应为 Inches() 等 Emu 值。"""
        if not items:
            return
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            self._style_paragraph(p, theme, "body")
            p.space_after = Pt(8)

    def _add_textbox(self, slide, text: str, left, top, width, height, theme: dict, role: str):
        """添加独立文本框，返回 shape。"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        self._style_paragraph(p, theme, role)
        return txBox

    def _add_page_number(self, slide, index: int, total: int, theme: dict) -> None:
        """在右下角添加页码。"""
        page_text = f"{index} / {total - 1}"
        txBox = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = page_text
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(10)
        p.font.color.rgb = theme["subtitle_color"]

    @staticmethod
    def _style_paragraph(para, theme: dict, role: str) -> None:
        """根据角色设置段落样式。"""
        if role == "title":
            para.font.size = Pt(30)
            para.font.bold = True
            para.font.color.rgb = theme["title_color"]
            para.font.name = theme["title_font"]
        elif role == "subtitle":
            para.font.size = Pt(20)
            para.font.color.rgb = theme["subtitle_color"]
            para.font.name = theme["body_font"]
        elif role == "body":
            para.font.size = Pt(18)
            para.font.color.rgb = theme["body_color"]
            para.font.name = theme["body_font"]

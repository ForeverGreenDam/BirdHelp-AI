"""RAG 摄取管道 — 文档下载→类型检测→LangChain 解析→清洗→切分→嵌入→入库。

对接 Java 后端文件服务，下载后进行向量化处理。
"""

from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from loguru import logger

from config import settings
from client.file import download as java_download
from rag.vector_store import add_documents
from utils.file import temp_file_path, ensure_temp_dir


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def _detect_loader(ext: str):
    """根据文件扩展名返回对应的 LangChain Loader 或自定义解析函数。"""
    if ext == ".pdf":
        from langchain_community.document_loaders import PyPDFLoader
        return PyPDFLoader
    if ext == ".docx":
        return _load_docx
    if ext == ".pptx":
        return _load_pptx
    if ext == ".txt":
        return _load_txt
    raise ValueError(f"不支持的文件格式: {ext}")


# ── 自定义 Loader（python-docx / python-pptx 直接提取文本） ──

def _load_docx(file_path: str) -> list[Document]:
    """从 .docx 提取所有段落文本。"""
    from docx import Document as DocxDocument
    doc = DocxDocument(file_path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [Document(page_content=text)]


def _load_pptx(file_path: str) -> list[Document]:
    """从 .pptx 提取所有幻灯片文本。"""
    from pptx import Presentation
    prs = Presentation(file_path)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
    return [Document(page_content="\n".join(parts))]


def _load_txt(file_path: str) -> list[Document]:
    """读取纯文本文件，自动尝试 UTF-8 / GBK 编码。"""
    for enc in ("utf-8", "gbk"):
        try:
            with open(file_path, "r", encoding=enc) as f:
                return [Document(page_content=f.read())]
        except UnicodeDecodeError:
            continue
    raise ValueError("无法解码文本文件")


# ── 切分器 ──

_splitter = None


def _get_splitter() -> RecursiveCharacterTextSplitter:
    """懒加载 RecursiveCharacterTextSplitter（全局单例）。"""
    global _splitter
    if _splitter is None:
        _splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
            separators=["\n\n", "\n", "。", ".", " ", ""],
        )
    return _splitter


# ── 主流程 ──

async def ingest_from_java(
    user_id: str,
    project_id: str,
    java_file_id: int,
    file_name: str = "unknown",
) -> dict:
    """从 Java 后端下载文件，解析→切分→嵌入→入库，返回摄取统计。

    Returns:
        {"file_id": int, "chunk_count": int, "file_name": str}
    """
    ext = Path(file_name).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件格式: {ext}，支持: {SUPPORTED_EXTENSIONS}")

    # 1. 下载到临时文件
    tmp_path = temp_file_path(ext)
    await java_download(java_file_id, save_path=str(tmp_path))
    logger.info(f"Downloaded file #{java_file_id} → {tmp_path}")

    try:
        # 2. 解析文档
        loader_cls = _detect_loader(ext)
        if ext == ".pdf":
            loader = loader_cls(str(tmp_path))
            raw_docs = loader.load()
        else:
            raw_docs = loader_cls(str(tmp_path))

        if not raw_docs:
            raise ValueError("文档解析结果为空，文件可能损坏或为空文件")

        # 3. 合并大文本 + 切分
        full_text = "\n\n".join(d.page_content for d in raw_docs)
        splitter = _get_splitter()
        chunks = splitter.create_documents(
            texts=[full_text],
            metadatas=[{
                "material_id": java_file_id,
                "user_id": user_id,
                "project_id": project_id,
                "file_name": file_name,
                "source": "java_upload",
            }],
        )
        # 为每个 chunk 追加 chunk_index
        for i, chunk in enumerate(chunks):
            chunk.metadata["chunk_index"] = i

        # 4. 嵌入 + 入库
        ids = add_documents(user_id, project_id, chunks)
        logger.info(f"Ingested file #{java_file_id} → project {project_id}: {len(ids)} chunks")

        return {
            "file_id": java_file_id,
            "chunk_count": len(ids),
            "file_name": file_name,
        }

    finally:
        # 5. 清理临时文件
        if tmp_path.exists():
            tmp_path.unlink()
